import streamlit as st
import re
import os
from openai import OpenAI
from PIL import Image
from pptx import Presentation
from streamlit_extras.stylable_container import stylable_container
import base64
from dotenv import load_dotenv
import time
import pandas as pd
import spacy
import io


import sys
print(">>> PYTHON VERSION =", sys.executable)

# Try to load spaCy for better NER-based anonymization
try:
    nlp = spacy.load("fr_core_news_md")
    SPACY_AVAILABLE = True
except (ImportError, OSError):
    SPACY_AVAILABLE = False

# Load environment variables (for OpenAI API key)
load_dotenv()

# Page setup
st.set_page_config(
    page_title="LinkedIn Post Generator",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for better appearance
st.markdown("""
<style>
    /* Main colors */
    :root {
        --main-color: #0A66C2;
        --accent-color: #2977C9;
        --bg-color: #f5f7f9;
        --text-color: #333;
        --light-gray: #e1e9f0;
    }
    
    /* Main container styling */
    .main {
        background-color: var(--bg-color);
        color: var(--text-color);
        font-family: "Inter", sans-serif;
    }
    
    /* Headers styling */
    h1, h2, h3 {
        color: var(--main-color);
        font-weight: 700;
    }
    
    /* Button styling */
    .stButton > button {
        background-color: var(--main-color);
        color: white;
        border-radius: 6px;
        padding: 0.5rem 1rem;
        border: none;
        font-weight: 600;
        transition: all 0.3s;
    }
    
    .stButton > button:hover {
        background-color: var(--accent-color);
        transform: translateY(-2px);
        box-shadow: 0 4px 8px rgba(10, 102, 194, 0.2);
    }
    
    /* Card styling */
    .css-1r6slb0 {
        border-radius: 10px;
        box-shadow: 0 4px 12px rgba(0, 0, 0, 0.05);
        background-color: white;
        padding: 1.5rem;
        margin-bottom: 1rem;
    }
    
    /* File uploader */
    .stFileUploader > div > button {
        background-color: white;
        color: var(--main-color);
        border: 2px solid var(--main-color);
    }
    
    /* Progress bar */
    .stProgress > div > div {
        background-color: var(--main-color);
    }
    
    /* Custom LinkedIn container */
    .linkedin-post-container {
        border: 1px solid #e1e9f0;
        border-radius: 8px;
        padding: 20px;
        background-color: white;
        box-shadow: 0 2px 6px rgba(0, 0, 0, 0.05);
    }
    
    .linkedin-header {
        display: flex;
        align-items: center;
        margin-bottom: 15px;
    }
    
    .linkedin-profile {
        width: 50px;
        height: 50px;
        border-radius: 50%;
        margin-right: 10px;
    }
    
    .linkedin-name {
        font-weight: 600;
        margin-bottom: 0;
    }
    
    .linkedin-headline {
        color: #666;
        font-size: 0.85rem;
        margin-top: 0;
    }
    
    .carousel-indicator {
        display: flex;
        justify-content: center;
        margin-top: 10px;
    }
    
    .carousel-dot {
        height: 8px;
        width: 8px;
        background-color: #bbb;
        border-radius: 50%;
        display: inline-block;
        margin: 0 4px;
    }
    
    .carousel-dot.active {
        background-color: var(--main-color);
    }
    
    /* Tab styling */
    .stTabs [data-baseweb="tab-list"] {
        gap: 8px;
    }
    
    .stTabs [data-baseweb="tab"] {
        background-color: white;
        border-radius: 6px 6px 0 0;
        padding: 8px 16px;
        border: 1px solid #e1e9f0;
        border-bottom: none;
    }
    
    .stTabs [aria-selected="true"] {
        background-color: var(--main-color) !important;
        color: white !important;
    }
    
    /* Step indicator */
    .step-container {
        display: flex;
        justify-content: space-between;
        margin-bottom: 2rem;
        position: relative;
    }
    
    .step-container:before {
        content: "";
        position: absolute;
        top: 15px;
        left: 0;
        right: 0;
        height: 2px;
        background-color: var(--light-gray);
        z-index: 1;
    }
    
    .step {
        display: flex;
        flex-direction: column;
        align-items: center;
        position: relative;
        z-index: 2;
    }
    
    .step-circle {
        width: 30px;
        height: 30px;
        border-radius: 50%;
        background-color: white;
        border: 2px solid var(--light-gray);
        display: flex;
        align-items: center;
        justify-content: center;
        margin-bottom: 8px;
        font-weight: 600;
        font-size: 14px;
    }
    
    .step.active .step-circle {
        background-color: var(--main-color);
        border-color: var(--main-color);
        color: white;
    }
    
    .step.completed .step-circle {
        background-color: var(--main-color);
        border-color: var(--main-color);
        color: white;
    }
    
    .step-label {
        font-size: 12px;
        color: #666;
        text-align: center;
    }
    
    .step.active .step-label {
        color: var(--main-color);
        font-weight: 600;
    }
</style>
""", unsafe_allow_html=True)

# Initialize session state variables if they don't exist
if 'active_step' not in st.session_state:
    st.session_state.active_step = 1
if 'extracted_text' not in st.session_state:
    st.session_state.extracted_text = ""
if 'anonymized_text' not in st.session_state:
    st.session_state.anonymized_text = ""
if 'linkedin_post' not in st.session_state:
    st.session_state.linkedin_post = ""
if 'carousel_content' not in st.session_state:
    st.session_state.carousel_content = []
if 'custom_entities' not in st.session_state:
    st.session_state.custom_entities = {}
if 'company_style' not in st.session_state:
    st.session_state.company_style = """
    Our LinkedIn posts typically:
    - Start with a thought-provoking question or bold statement
    - Use professional but conversational language
    - Include specific results and metrics when possible
    - End with a clear call to action
    - Use 3-5 relevant hashtags, including #YourCompanyName
    """

# Function to display progress tracker
def display_progress_tracker():
    # Updated to skip the slide selection step
    steps = [
        "Upload Presentation", 
        "Extract Text", 
        "Anonymize Text", 
        "Generate Post", 
        "Review & Edit"
    ]
    
    st.markdown('<div class="step-container">', unsafe_allow_html=True)
    
    for i, step in enumerate(steps, 1):
        if i < st.session_state.active_step:
            status = "completed"
        elif i == st.session_state.active_step:
            status = "active"
        else:
            status = ""
        
        st.markdown(f'''
        <div class="step {status}">
            <div class="step-circle">{i}</div>
            <div class="step-label">{step}</div>
        </div>
        ''', unsafe_allow_html=True)
    
    st.markdown('</div>', unsafe_allow_html=True)

# Function to extract text from PowerPoint slides
def extract_text_from_presentation(presentation):
    extracted_texts = []
    
    for slide in presentation.slides:
        slide_text = ""
        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        
        # Only add non-empty slides
        if slide_text.strip():
            extracted_texts.append(slide_text.strip())
    
    return "\n\n".join(extracted_texts)

# Enhanced anonymization function using spaCy if available
def anonymize_text(text, custom_entities=None):
    # First ensure spaCy is properly installed and loaded
    try:
        import spacy
        try:
            nlp = spacy.load("fr_core_news_md")
            SPACY_AVAILABLE = True
        except OSError:
            st.warning("French language model not found. Attempting to download automatically...")
            try:
                from spacy.cli import download
                download("fr_core_news_md")
                nlp = spacy.load("fr_core_news_md")
                SPACY_AVAILABLE = True
                st.success("French language model downloaded successfully! Enhanced anonymization enabled.")
            except Exception as e:
                st.error(f"Could not automatically download the French language model: {str(e)}")
                st.info("Falling back to basic anonymization. Check the 'How to enable enhanced anonymization' section for manual installation instructions.")
                SPACY_AVAILABLE = False
    except ImportError:
        st.warning("spaCy library not found. Using basic anonymization only.")
        SPACY_AVAILABLE = False
    
    anonymized = text
    
    # First apply spaCy NER if available
    if SPACY_AVAILABLE:
        doc = nlp(anonymized)
        
        # Create a list of replacements to avoid modifying the string while iterating
        replacements = []
        
        # Define entity mappings (expanded)
        entity_mappings = {
            'PERSON': '[PERSONNE]',
            'ORG': '[ORGANISATION]',
            'GPE': '[LIEU]',
            'LOC': '[LOCALISATION]',
            'PRODUCT': '[PRODUIT]',
            'MONEY': '[MONTANT]',
            'CARDINAL': '[NOMBRE]',
            'DATE': '[DATE]',
            'TIME': '[HEURE]',
            'PERCENT': '[POURCENTAGE]',
            'QUANTITY': '[QUANTITE]',
            'NORP': '[GROUPE]',  # Nationalities, religious groups
            'FAC': '[INSTALLATION]',  # Buildings, airports, highways
            'EVENT': '[EVENEMENT]',
            'WORK_OF_ART': '[OEUVRE]',
            'LAW': '[LOI]',
            'LANGUAGE': '[LANGUE]',
            'ORDINAL': '[ORDINAL]'
        }
        
        # Process entities found by spaCy
        for ent in doc.ents:
            if ent.label_ in entity_mappings:
                replacements.append((ent.start_char, ent.end_char, entity_mappings[ent.label_]))
        
        # Apply replacements in reverse order to maintain correct indices
        replacements.sort(key=lambda x: x[0], reverse=True)
        for start, end, replacement in replacements:
            anonymized = anonymized[:start] + replacement + anonymized[end:]
    
    # Custom entity replacement (with more robust approach)
    if custom_entities:
        for entity, replacement in custom_entities.items():
            if entity and entity.strip():
                # Use regex with word boundaries for more precise replacement
                pattern = r'\b' + re.escape(entity) + r'\b'
                anonymized = re.sub(pattern, replacement, anonymized, flags=re.IGNORECASE)
    
    # Enhanced regex patterns for better coverage
    patterns = [
        # Emails with broader pattern
        (r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b', '[EMAIL]'),
        
        # Various phone number formats (international, French, etc.)
        (r'\b(?:\+\d{1,3}[\s.-]?)?\(?\d{2,4}\)?[\s.-]?\d{2,3}[\s.-]?\d{2,4}(?:[\s.-]?\d{2,4})?\b', '[TÉLÉPHONE]'),
        
        # URLs (more comprehensive)
        (r'https?://[^\s<>"]+|www\.[^\s<>"]+', '[URL]'),
        
        # Dates in various formats
        (r'\b(?:\d{1,2}[/.-]\d{1,2}[/.-]\d{2,4}|\d{4}[/.-]\d{1,2}[/.-]\d{1,2})\b', '[DATE]'),
        (r'\b(?:janvier|février|mars|avril|mai|juin|juillet|août|septembre|octobre|novembre|décembre)\s+\d{4}\b', '[DATE]'),
        
        # Currency amounts (various formats)
        (r'\b\d+(?:[\s,.]\d+)*(?:\s*€|\s*EUR|\s*\$|\s*USD|\s*MAD|\s*dirhams?)\b', '[MONTANT]'),
        
        # Percentages
        (r'\b\d+(?:[,.]\d+)?%\b', '[POURCENTAGE]'),
        
        # Postal codes (various formats)
        (r'\b\d{5}\b|\b[A-Z]\d[A-Z]\s?\d[A-Z]\d\b', '[CODE POSTAL]'),
        
        # Company registration numbers
        (r'\b\d{9}(?:\d{5})?\b', '[SIRET/SIREN]'),
        
        # Addresses (more comprehensive)
        (r'\b\d+\s+(?:rue|avenue|boulevard|impasse|allée|place|chemin|voie)\b.{3,50}?(?:,|\.|$)', '[ADRESSE]'),
        
        # Credit card numbers
        (r'\b(?:\d{4}[\s-]?){4}\b', '[CARTE DE CRÉDIT]'),
        
        # IP addresses
        (r'\b(?:\d{1,3}\.){3}\d{1,3}\b', '[ADRESSE IP]'),
        
        # Social security/identity numbers
        (r'\b\d\s\d{2}\s\d{2}\s\d{2}\s\d{3}\s\d{3}\s\d{2}\b', '[NUMÉRO DE SÉCURITÉ SOCIALE]')
    ]
    
    # Apply all patterns
    for pattern, replacement in patterns:
        anonymized = re.sub(pattern, replacement, anonymized, flags=re.IGNORECASE)
    
    # Look for price ranges
    anonymized = re.sub(r'\b\d+\s*(?:à|to|-)\s*\d+\s*(?:€|EUR|\$|USD|MAD|dirhams?)\b', '[FOURCHETTE PRIX]', anonymized, flags=re.IGNORECASE)
    
    # Additional brands/products (expanded list)
    brands = [
        'TONIK', 'TOBIGO', 'LAGO POKER', 'LOACKER', 'MANNER', 'TAGGER', 
        'LAGO PLAISIR', 'BONO', 'KITKAT', 'COCA-COLA', 'PEPSI', 'NESTLÉ',
        'DANONE', 'APPLE', 'SAMSUNG', 'MICROSOFT', 'AMAZON', 'GOOGLE',
        'FACEBOOK', 'INSTAGRAM', 'TWITTER', 'LINKEDIN', 'YOUTUBE',
        'ADIDAS', 'NIKE', 'PUMA', 'REEBOK', 'ZARA', 'H&M'
    ]
    
    for brand in brands:
        pattern = r'\b' + re.escape(brand) + r'\b'
        anonymized = re.sub(pattern, '[MARQUE]', anonymized, flags=re.IGNORECASE)
    
    # More cities (expanded)
    cities = [
        'Casablanca', 'Agadir', 'Marrakech', 'Rabat', 'Fès', 'Tanger', 'Meknès',
        'Oujda', 'Tétouan', 'Nador', 'Kénitra', 'El Jadida', 'Béni Mellal',
        'Mohammedia', 'Essaouira', 'Ouarzazate', 'Paris', 'Lyon', 'Marseille',
        'Toulouse', 'Nice', 'Bordeaux', 'Lille', 'New York', 'London', 'Madrid'
    ]
    
    for city in cities:
        pattern = r'\b' + re.escape(city) + r'\b'
        anonymized = re.sub(pattern, '[VILLE]', anonymized, flags=re.IGNORECASE)
    
    # Masquer les citations en bloc (verbatim utilisateur)
    anonymized = re.sub(r'«[^»]{5,500}»', '[CITATION UTILISATEUR]', anonymized)
    anonymized = re.sub(r'"[^"]{5,500}"', '[CITATION UTILISATEUR]', anonymized)
    anonymized = re.sub(r"'[^']{5,500}'", '[CITATION UTILISATEUR]', anonymized)
    
    return anonymized
# Function to generate LinkedIn post via OpenAI
def generate_linkedin_post(anonymized_text, post_type="case_study", company_style=""):
    try:
        client = OpenAI(api_key=st.secrets.get("openai_api_key", os.getenv("OPENAI_API_KEY")))
        
        # Different prompt templates based on post type
        prompts = {
            "case_study": f"""
            Based on the following anonymized content from a case study presentation, create:
            
            1. A compelling LinkedIn post text (300-400 words) that highlights the key achievements, 
               challenges overcome, and business impact of this project
            2. Content for 3-5 carousel slides that would accompany this LinkedIn post
            
            Follow our company style guidelines:
            {company_style}
            
            Content from presentation:
            {anonymized_text}
            """,
            
            "product_launch": f"""
            Based on the following anonymized content from a product presentation, create:
            
            1. An exciting LinkedIn post text (300-400 words) that builds anticipation for our new 
               product launch, highlighting key features and benefits
            2. Content for 3-5 carousel slides that would showcase the product's unique selling points
            
            Follow our company style guidelines:
            {company_style}
            
            Content from presentation:
            {anonymized_text}
            """,
            
            "thought_leadership": f"""
            Based on the following anonymized content from a presentation, create:
            
            1. A thought-provoking LinkedIn post text (300-400 words) that positions our company 
               as an industry thought leader with valuable insights
            2. Content for 3-5 carousel slides that would outline key industry trends or insights
            
            Follow our company style guidelines:
            {company_style}
            
            Content from presentation:
            {anonymized_text}
            """
        }
        
        # Select the appropriate prompt
        selected_prompt = prompts.get(post_type, prompts["case_study"])
        
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "Vous êtes un créateur de contenu professionnel spécialisé dans le contenu LinkedIn. Votre tâche est de créer du contenu concis, percutant et professionnel pour des posts et des carrousels LinkedIn qui respectent le style de l'entreprise."},
                {"role": "user", "content": selected_prompt}
            ],
            max_tokens=1500,
            temperature=0.7
        )
        
        result = response.choices[0].message.content.strip()
        
        # Split the response into post text and carousel content
        post_text_match = re.search(r'POST TEXT:(.*?)CAROUSEL SLIDES:', result, re.DOTALL)
        carousel_match = re.search(r'CAROUSEL SLIDES:(.*)', result, re.DOTALL)
        
        post_text = post_text_match.group(1).strip() if post_text_match else ""
        
        carousel_content = []
        if carousel_match:
            carousel_text = carousel_match.group(1).strip()
            slide_pattern = r'Slide \d+: (.*?)\n(.*?)(?=Slide \d+:|$)'
            slides = re.findall(slide_pattern, carousel_text, re.DOTALL)
            
            for title, content in slides:
                carousel_content.append({
                    "title": title.strip(),
                    "content": content.strip()
                })
        
        return post_text, carousel_content
        
    except Exception as e:
        st.error(f"Error generating LinkedIn post: {str(e)}")
        return "", []

# Function to encode image to base64
def encode_image(image_bytes):
    return base64.b64encode(image_bytes).decode('utf-8')

# Function to create a LinkedIn post preview
def linkedin_post_preview(post_text):
    html_text = post_text.replace('\n', '<br>')

    st.markdown('<div class="linkedin-post-container">', unsafe_allow_html=True)
    
    # Profile header
    st.markdown('''
    <div class="linkedin-header">
        <img src="https://cdn-icons-png.flaticon.com/512/149/149071.png" class="linkedin-profile">
        <div>
            <p class="linkedin-name">Your Name</p>
            <p class="linkedin-headline">Your Professional Title</p>
        </div>
    </div>
    ''', unsafe_allow_html=True)
    
        # Post content
    st.markdown(f"<p>{html_text}</p>", unsafe_allow_html=True)
    
    # Carousel indicator
    st.markdown('<div class="carousel-indicator">', unsafe_allow_html=True)
    for i in range(len(st.session_state.carousel_content)):
        if i == 0:
            st.markdown('<span class="carousel-dot active"></span>', unsafe_allow_html=True)
        else:
            st.markdown('<span class="carousel-dot"></span>', unsafe_allow_html=True)
    st.markdown('</div>', unsafe_allow_html=True)
    
    st.markdown('</div>', unsafe_allow_html=True)

# Function to display carousel slide preview
def carousel_slide_preview(slide):
    with st.container():
        content_html = slide['content'].replace('\n', '<br>')
        st.markdown(
            f"""
            <div style="border: 1px solid #e1e9f0; border-radius: 8px; padding: 20px; background-color: white; text-align: center;">
                <h3 style="color: #0A66C2;">{slide['title']}</h3>
                <p>{content_html}</p>
            </div>
            """,
            unsafe_allow_html=True
        )

# Function to move to next step
def next_step():
    st.session_state.active_step += 1

# Function to go back to previous step
def prev_step():
    if st.session_state.active_step > 1:
        st.session_state.active_step -= 1

# Function to create downloadable carousel slides
def create_carousel_slides():
    slides_data = []
    
    for i, slide in enumerate(st.session_state.carousel_content):
        # Create a PIL image with text
        img_width, img_height = 1200, 900
        img = Image.new('RGB', (img_width, img_height), color=(255, 255, 255))
        
        try:
            # Draw text on the image
            from PIL import ImageDraw, ImageFont
            draw = ImageDraw.Draw(img)
            
            # Try to load a nice font, fallback to default
            try:
                title_font = ImageFont.truetype(r"Arial.ttf", 60)
                content_font = ImageFont.truetype(r"Arial.ttf", 40)
            except IOError:
                # Use default font if Arial is not available
                title_font = ImageFont.load_default()
                content_font = ImageFont.load_default()
            
            # Draw title
            title = slide['title']
            title_width = draw.textlength(title, font=title_font)
            draw.text(((img_width - title_width) // 2, 100), title, fill=(10, 102, 194), font=title_font)
            
            # Draw content (simple multi-line text)
            content = slide['content']
            content_lines = content.split('\n')
            y_position = 250
            for line in content_lines:
                line_width = draw.textlength(line, font=content_font)
                draw.text(((img_width - line_width) // 2, y_position), line, fill=(0, 0, 0), font=content_font)
                y_position += 50
            
            # Convert to bytes for download
            img_byte_arr = io.BytesIO()
            img.save(img_byte_arr, format='PNG')
            img_bytes = img_byte_arr.getvalue()
            
            slides_data.append((f"slide_{i+1}.png", img_bytes))
        
        except Exception as e:
            st.error(f"Error creating slide image: {str(e)}")
    
    return slides_data

# Main app header
st.title("LinkedIn Post Generator from PowerPoint")
st.markdown("Transform your presentation into engaging LinkedIn content with just a few clicks.")

# Display progress tracker
display_progress_tracker()

# Step 1: Upload PowerPoint presentation
if st.session_state.active_step == 1:
    with st.container():
        st.header("Step 1: Upload Your PowerPoint Presentation")
        
        uploaded_file = st.file_uploader("Choose a PowerPoint file", type=["ppt", "pptx"])
        
        if uploaded_file is not None:
            # Save the uploaded file to a temporary file
            with open("temp.pptx", "wb") as f:
                f.write(uploaded_file.getvalue())
            
            # Load the presentation
            presentation = Presentation("temp.pptx")
            
            # Store presentation in session state
            st.session_state.presentation = presentation
            
            # Count slides with content
            content_slides = 0
            for slide in presentation.slides:
                has_content = False
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        has_content = True
                        break
                if has_content:
                    content_slides += 1
            
            st.success(f"Successfully loaded presentation with {len(presentation.slides)} slides ({content_slides} with text content)!")
            
            if st.button("Continue to Extract Text"):
                next_step()

# Step 2: Extract Text (skipping slide selection)
elif st.session_state.active_step == 2:
    st.header("Step 2: Extract Text from Presentation")
    
    if hasattr(st.session_state, 'presentation'):
        with st.spinner("Extracting text from presentation..."):
            # Extract text from all slides
            extracted_text = extract_text_from_presentation(st.session_state.presentation)
            st.session_state.extracted_text = extracted_text
            
            st.success("Text extracted successfully!")
            
            st.subheader("Extracted Text:")
            
            # Allow editing the extracted text
            edited_text = st.text_area("You can edit the extracted text if needed:", 
                                     value=extracted_text, 
                                     height=300, 
                                     key="extracted_text_display")
            
            # Update session state if text was edited
            if edited_text != st.session_state.extracted_text:
                st.session_state.extracted_text = edited_text
            
            col1, col2 = st.columns([1, 5])
            with col1:
                if st.button("← Back"):
                    prev_step()
            with col2:
                if st.button("Continue to Anonymization"):
                    next_step()
    else:
        st.warning("No presentation loaded. Please upload a presentation first.")
        if st.button("← Back to Upload"):
            prev_step()

# Step 3: Anonymize Text
elif st.session_state.active_step == 3:
    st.header("Step 3: Anonymize Text")
    
    if st.session_state.extracted_text:
        st.info("Add any specific terms you want to anonymize or replace. These will be applied in addition to the automatic anonymization.")
        
        # Custom entities for anonymization
        with st.expander("Custom Anonymization Rules", expanded=True):
            col1, col2 = st.columns(2)
            
            with col1:
                custom_entity = st.text_input("Enter term to anonymize:")
            with col2:
                replacement = st.text_input("Replace with:")
            
            if st.button("Add Rule"):
                if custom_entity and custom_entity.strip():
                    st.session_state.custom_entities[custom_entity] = replacement
                    st.success(f"Added rule: '{custom_entity}' → '{replacement}'")
            
            # Display current custom entities
            if st.session_state.custom_entities:
                st.subheader("Current Custom Rules:")
                for entity, repl in st.session_state.custom_entities.items():
                    col1, col2, col3 = st.columns([5, 5, 1])
                    with col1:
                        st.text(entity)
                    with col2:
                        st.text(f"→ {repl}")
                    with col3:
                        if st.button("🗑️", key=f"delete_{entity}"):
                            del st.session_state.custom_entities[entity]
                            st.rerun()
        
        # Apply anonymization with both automatic and custom rules
        anonymized_text = anonymize_text(
            st.session_state.extracted_text, 
            st.session_state.custom_entities
        )
        
        st.session_state.anonymized_text = anonymized_text
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.subheader("Original Text:")
            st.text_area("Original content", value=st.session_state.extracted_text, height=300, key="original_text_display")
        
        with col2:
            st.subheader("Anonymized Text:")
            edited_anon_text = st.text_area("Anonymized content (you can edit this)", value=anonymized_text, height=300, key="anonymized_text_display")
            
            # Update if user manually edited
            if edited_anon_text != anonymized_text:
                st.session_state.anonymized_text = edited_anon_text
        
        st.info("Note: If spaCy is installed, entity recognition will be used for better anonymization. Otherwise, only regex-based rules will be applied.")
        
        col1, col2 = st.columns([1, 5])
        with col1:
            if st.button("← Back"):
                prev_step()
        with col2:
            if st.button("Generate LinkedIn Post"):
                next_step()
    else:
        st.warning("No extracted text available. Please go back to extract text.")
        if st.button("← Back to Text Extraction"):
            prev_step()

# Step 4: Generate LinkedIn Post
elif st.session_state.active_step == 4:
    st.header("Step 4: Generate LinkedIn Post")
    
    if st.session_state.anonymized_text:
        # Post type selection
        post_type = st.radio(
            "Select the type of LinkedIn post to generate:",
            ["case_study", "product_launch", "thought_leadership"],
            horizontal=True
        )
        
        # Company style customization
        with st.expander("Customize Company Style"):
            st.session_state.company_style = st.text_area(
                "Customize your company's LinkedIn post style:",
                value=st.session_state.company_style,
                height=150
            )
        
        # Generate button
        if st.button("Generate LinkedIn Post"):
            with st.spinner("Generating LinkedIn post..."):
                # Add a progress bar for visual effect
                progress_bar = st.progress(0)
                for i in range(100):
                    time.sleep(0.02)  # Simulating work being done
                    progress_bar.progress(i + 1)
                
                # Generate LinkedIn post using OpenAI
                post_text, carousel_content = generate_linkedin_post(
                    st.session_state.anonymized_text,
                    post_type,
                    st.session_state.company_style
                )
                
                st.session_state.linkedin_post = post_text
                st.session_state.carousel_content = carousel_content
        
        # Display the generated content
        if st.session_state.linkedin_post:
            st.success("LinkedIn post generated successfully!")
            
            st.subheader("Generated LinkedIn Post:")
            post_text = st.text_area("Post text", value=st.session_state.linkedin_post, height=200, key="post_text_display")
            
            # Update the session state with any edits
            st.session_state.linkedin_post = post_text
            
            st.subheader("Generated Carousel Slides:")
            
            # Display carousel slides
            for i, slide in enumerate(st.session_state.carousel_content):
                with st.expander(f"Slide {i+1}: {slide['title']}", expanded=True):
                    title = st.text_input("Title", value=slide['title'], key=f"slide_title_{i}")
                    content = st.text_area("Content", value=slide['content'], height=100, key=f"slide_content_{i}")
                    
                    # Update the session state with any edits
                    st.session_state.carousel_content[i]['title'] = title
                    st.session_state.carousel_content[i]['content'] = content
            
            col1, col2 = st.columns([1, 5])
            with col1:
                if st.button("← Back"):
                    prev_step()
            with col2:
                if st.button("Continue to Preview"):
                    next_step()
        else:
            st.info("Click 'Generate LinkedIn Post' to create content based on your anonymized text.")
            
            if st.button("← Back to Anonymization"):
                prev_step()
    else:
        st.warning("No anonymized text available. Please go back to anonymize text.")
        if st.button("← Back to Anonymization"):
            prev_step()

# Step 5: Review and Edit
elif st.session_state.active_step == 5:
    st.header("Step 5: Review and Edit Final Content")
    
    if st.session_state.linkedin_post and st.session_state.carousel_content:
        # Display tabs for different views
        tab1, tab2, tab3 = st.tabs(["LinkedIn Post Preview", "Carousel Content", "Export Options"])
        
        with tab1:
            st.subheader("LinkedIn Post Preview")
            linkedin_post_preview(st.session_state.linkedin_post)
            
            # Option to edit post text
            edited_post = st.text_area("Edit Post Text", value=st.session_state.linkedin_post, height=200)
            if edited_post != st.session_state.linkedin_post:
                st.session_state.linkedin_post = edited_post
                st.rerun()
        
        with tab2:
            st.subheader("Carousel Slides")
            
            # Display carousel slides with edit options
            for i, slide in enumerate(st.session_state.carousel_content):
                st.markdown(f"### Slide {i+1}")
                carousel_slide_preview(slide)
                
                col1, col2 = st.columns(2)
                with col1:
                    edited_title = st.text_input("Edit Title", value=slide['title'], key=f"edit_title_{i}")
                with col2:
                    edited_content = st.text_area("Edit Content", value=slide['content'], height=100, key=f"edit_content_{i}")
                
                if edited_title != slide['title'] or edited_content != slide['content']:
                    st.session_state.carousel_content[i]['title'] = edited_title
                    st.session_state.carousel_content[i]['content'] = edited_content
                    st.rerun()
        
        with tab3:
            st.subheader("Export Options")
            
            # Post text export
            col1, col2 = st.columns(2)
            with col1:
                st.download_button(
                    "Download Post Text",
                    st.session_state.linkedin_post,
                    file_name="linkedin_post.txt",
                    mime="text/plain"
                )
            
            # Generate and provide carousel slides as images
            with col2:
                slides_data = create_carousel_slides()
                
                for slide_name, slide_bytes in slides_data:
                    st.download_button(
                        f"Download {slide_name}",
                        slide_bytes,
                        file_name=slide_name,
                        mime="image/png"
                    )
            
            # Export all as a zip
            if len(slides_data) > 0:
                try:
                    import zipfile
                    from io import BytesIO
                    
                    # Create in-memory zip file
                    zip_buffer = BytesIO()
                    with zipfile.ZipFile(zip_buffer, 'a', zipfile.ZIP_DEFLATED) as zip_file:
                        # Add post text
                        zip_file.writestr("linkedin_post.txt", st.session_state.linkedin_post)
                        
                        # Add slide images
                        for slide_name, slide_bytes in slides_data:
                            zip_file.writestr(slide_name, slide_bytes)
                    
                    # Download complete package
                    st.download_button(
                        "Download Complete Package (ZIP)",
                        zip_buffer.getvalue(),
                        file_name="linkedin_post_package.zip",
                        mime="application/zip"
                    )
                except Exception as e:
                    st.error(f"Error creating ZIP package: {str(e)}")
        
        # Navigation buttons
        col1, col2 = st.columns([1, 5])
        with col1:
            if st.button("← Back"):
                prev_step()
        with col2:
            if st.button("Start Over"):
                for key in ['active_step', 'extracted_text', 'anonymized_text', 
                           'linkedin_post', 'carousel_content']:
                    if key in st.session_state:
                        if key == 'active_step':
                            st.session_state[key] = 1
                        elif key == 'custom_entities' or key == 'company_style':
                            # Keep these settings
                            pass
                        else:
                            st.session_state[key] = []
                st.rerun()
    else:
        st.warning("No LinkedIn post content available. Please go back to generate content.")
        if st.button("← Back to Generate LinkedIn Post"):
            prev_step()

# Add sidebar for app information and tips
with st.sidebar:
    st.title("About")
    st.markdown("""
    ## LinkedIn Post Generator
    
    This app helps you transform your PowerPoint presentations into engaging LinkedIn content. Perfect for:
    
    - Marketing teams
    - Sales professionals
    - Content creators
    - Business development
    
    ### How it works
    
    1. Upload your PowerPoint presentation
    2. Extract text from all slides
    3. Anonymize sensitive information
    4. Generate professional LinkedIn content
    5. Review, edit and export
    
    ### Tips for best results
    
    - Use presentations with clear, well-structured content
    - Review the anonymized text carefully
    - Choose the post type that best matches your content
    - Customize the company style to match your brand voice
    """)
    
    st.markdown("---")
    st.markdown("© 2025 LinkedIn Post Generator | Made with ❤️ and Streamlit")

# Footer
st.markdown("---")
st.markdown("© 2025 LinkedIn Post Generator | Made with ❤️ and Streamlit", unsafe_allow_html=True)
